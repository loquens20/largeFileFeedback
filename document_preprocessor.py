"""
파일 형식별 전처리기
- DOCX, PDF, PPTX 등 문서에서 텍스트와 이미지 추출
- 순서 유지
- 이미지 base64 인코딩
"""

import io
import base64
from pathlib import Path
from typing import List, Dict, Any, Tuple
from dataclasses import dataclass
from PIL import Image


@dataclass
class ExtractedContent:
    """추출된 콘텐츠"""
    position: int
    content_type: str  # 'text', 'image'
    data: Any  # 텍스트 또는 이미지 데이터


class DocumentPreprocessor:
    """문서 전처리 클래스"""

    def __init__(self, max_image_size: Tuple[int, int] = (2048, 2048)):
        self.max_image_size = max_image_size

    def resize_image(self, image: Image.Image) -> Image.Image:
        """이미지 크기 조정 (비용 절감)"""
        if image.size[0] > self.max_image_size[0] or image.size[1] > self.max_image_size[1]:
            image.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)
        return image

    def image_to_base64(self, image: Image.Image, format: str = 'PNG') -> str:
        """이미지를 base64로 인코딩"""
        buffered = io.BytesIO()
        image.save(buffered, format=format)
        return base64.b64encode(buffered.getvalue()).decode('utf-8')

    def extract_from_docx(self, file_path: str) -> List[ExtractedContent]:
        """DOCX 파일에서 텍스트와 이미지 추출"""
        try:
            from docx import Document
            from docx.oxml.text.paragraph import CT_P
            from docx.oxml.table import CT_Tbl
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
        except ImportError:
            raise ImportError("python-docx가 필요합니다: pip install python-docx")

        doc = Document(file_path)
        contents = []
        position = 0

        # 문서의 모든 요소를 순서대로 처리
        for element in doc.element.body:
            if isinstance(element, CT_P):
                # 단락
                para = Paragraph(element, doc)
                text = para.text.strip()

                if text:
                    contents.append(ExtractedContent(
                        position=position,
                        content_type='text',
                        data=text
                    ))
                    position += 1

                # 단락 내 이미지
                for run in para.runs:
                    if run._element.xpath('.//a:blip'):
                        for rel in run.part.rels.values():
                            if "image" in rel.target_ref:
                                image_data = rel.target_part.blob
                                image = Image.open(io.BytesIO(image_data))
                                image = self.resize_image(image)

                                contents.append(ExtractedContent(
                                    position=position,
                                    content_type='image',
                                    data=self.image_to_base64(image)
                                ))
                                position += 1

            elif isinstance(element, CT_Tbl):
                # 표
                table = Table(element, doc)
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(' | '.join(row_text))

                if table_text:
                    contents.append(ExtractedContent(
                        position=position,
                        content_type='text',
                        data='[표]\n' + '\n'.join(table_text)
                    ))
                    position += 1

        return contents

    def extract_from_pdf(self, file_path: str) -> List[ExtractedContent]:
        """PDF 파일에서 텍스트와 이미지 추출"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF가 필요합니다: pip install PyMuPDF")

        doc = fitz.open(file_path)
        contents = []
        position = 0

        for page_num in range(len(doc)):
            page = doc[page_num]

            # 페이지 번호 추가
            contents.append(ExtractedContent(
                position=position,
                content_type='text',
                data=f'\n--- 페이지 {page_num + 1} ---\n'
            ))
            position += 1

            # 텍스트 추출
            text = page.get_text()
            if text.strip():
                contents.append(ExtractedContent(
                    position=position,
                    content_type='text',
                    data=text
                ))
                position += 1

            # 이미지 추출
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]

                image = Image.open(io.BytesIO(image_bytes))
                image = self.resize_image(image)

                contents.append(ExtractedContent(
                    position=position,
                    content_type='image',
                    data=self.image_to_base64(image)
                ))
                position += 1

        doc.close()
        return contents

    def extract_from_pptx(self, file_path: str) -> List[ExtractedContent]:
        """PPTX 파일에서 텍스트와 이미지 추출"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx가 필요합니다: pip install python-pptx")

        prs = Presentation(file_path)
        contents = []
        position = 0

        for slide_num, slide in enumerate(prs.slides):
            # 슬라이드 번호
            contents.append(ExtractedContent(
                position=position,
                content_type='text',
                data=f'\n--- 슬라이드 {slide_num + 1} ---\n'
            ))
            position += 1

            # 도형에서 텍스트 추출
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    contents.append(ExtractedContent(
                        position=position,
                        content_type='text',
                        data=shape.text
                    ))
                    position += 1

                # 이미지 추출
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    image_bytes = image.blob

                    img = Image.open(io.BytesIO(image_bytes))
                    img = self.resize_image(img)

                    contents.append(ExtractedContent(
                        position=position,
                        content_type='image',
                        data=self.image_to_base64(img)
                    ))
                    position += 1

        return contents

    def extract_from_xlsx(self, file_path: str) -> List[ExtractedContent]:
        """XLSX 파일에서 텍스트 추출"""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl이 필요합니다: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path, data_only=True)
        contents = []
        position = 0

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]

            contents.append(ExtractedContent(
                position=position,
                content_type='text',
                data=f'\n--- 시트: {sheet_name} ---\n'
            ))
            position += 1

            # 데이터 추출
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                if any(row_data):  # 빈 행 제외
                    rows.append(' | '.join(row_data))

            if rows:
                contents.append(ExtractedContent(
                    position=position,
                    content_type='text',
                    data='\n'.join(rows)
                ))
                position += 1

        return contents

    def extract_content(self, file_path: str) -> List[ExtractedContent]:
        """파일 형식에 따라 적절한 추출 메서드 호출"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        extractors = {
            '.docx': self.extract_from_docx,
            '.pdf': self.extract_from_pdf,
            '.pptx': self.extract_from_pptx,
            '.xlsx': self.extract_from_xlsx,
            '.txt': self._extract_from_text,
        }

        if extension not in extractors:
            raise ValueError(f"지원하지 않는 파일 형식: {extension}")

        print(f"📄 파일 추출 중: {file_path.name} ({extension})")
        contents = extractors[extension](str(file_path))

        # 통계
        text_count = sum(1 for c in contents if c.content_type == 'text')
        image_count = sum(1 for c in contents if c.content_type == 'image')
        print(f"✓ 추출 완료: 텍스트 {text_count}개, 이미지 {image_count}개")

        return contents

    def _extract_from_text(self, file_path: str) -> List[ExtractedContent]:
        """일반 텍스트 파일 처리"""
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()

        return [ExtractedContent(
            position=0,
            content_type='text',
            data=text
        )]


# 사용 예제
if __name__ == "__main__":
    preprocessor = DocumentPreprocessor()

    print("전처리기 준비 완료")
    print("\n지원 형식:")
    print("- DOCX: python-docx 필요")
    print("- PDF: PyMuPDF 필요")
    print("- PPTX: python-pptx 필요")
    print("- XLSX: openpyxl 필요")
    print("- TXT: 기본 지원")
